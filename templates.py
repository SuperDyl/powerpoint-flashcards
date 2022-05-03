"""
Templates used to create each PowerPoint slide
"""

from pptxtemplate import SlideTemplate, TextboxTemplate, Position, Dimension

from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.shapes.autoshape import Shape as Textbox
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT

from typing import Callable
from copy import copy


class TextboxBaseTemplate:
    """
    Base formatting for all text boxes in PowerPoint Flashcards

    """
    def __init__(self, template: TextboxTemplate):
        self.template = copy(template)
        self.font_size = Pt(72)
        self.font_color = RGBColor(255, 0, 0)
        self.font_color_mode = "rgb"  # FONT_COLOR_MODE["RGB"]
        self.name_bottom_padding = Inches(1.58)

    def add_base_textbox(self, slide: Slide, presentation: Presentation) -> Textbox:
        """Add to slide and return a textbox with default formatting"""
        self.reset_textbox_position(presentation)

        textbox = self.template.add_textbox(slide)

        frame = textbox.text_frame
        frame.auto_size = MSO_AUTO_SIZE.NONE
        frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

        font = frame.paragraphs[0].font
        font.size = self.font_size
        setattr(font.color, self.font_color_mode, self.font_color)

        return textbox

    def reset_textbox_position(self, presentation: Presentation) -> None:
        textbox_height = self.name_bottom_padding
        self.template.top = presentation.slide_height - textbox_height
        self.template.width = presentation.slide_width
        self.template.height = textbox_height


def build_professor_textbox_func(template: TextboxTemplate)\
        -> Callable[[Slide, Presentation, str], Textbox]:
    """Return a function to format, add, and return a textbox formatted for a professor's name"""
    base_template = TextboxBaseTemplate(template)

    def add_professor_name_textbox(slide: Slide, presentation: Presentation, professor_name: str)\
            -> Textbox:
        """Add to slide and then return a default formatted textbox"""
        textbox = base_template.add_base_textbox(slide, presentation)
        textbox.text_frame.paragraphs[0].add_run().text = professor_name
        return textbox

    return add_professor_name_textbox


def build_room_textbox_func(template: TextboxTemplate) -> Callable[[Slide, Presentation, str], Textbox]:
    """Return a function to format, add, and return a textbox formatted for a room number"""
    base_template = TextboxBaseTemplate(template)

    def add_room_number_textbox(slide: Slide, presentation: Presentation, room_number: str) -> Textbox:
        """Add to slide and then return a default formatted textbox"""
        textbox = base_template.add_base_textbox(slide, presentation)
        textbox.text_frame.paragraphs[0].add_run().text = room_number
        return textbox

    return add_room_number_textbox


class BaseSlideTemplate:
    """
    Add and return a new slide formatted for a room
    """
    def __init__(self, slide_template: SlideTemplate, position: Position = Position(left=Inches(0), top=Inches(0)),
                 dimension: Dimension = Dimension(width=Inches(10), height=Inches(1))):
        self.slide_template = slide_template
        self.name_bottom_padding = Inches(1.58)

        self.professor_base_template = TextboxTemplate(position, dimension)
        self.professor_name_template = build_professor_textbox_func(self.professor_base_template)

    def add_base_slide(self, presentation: Presentation, professor_name: str) -> Slide:
        """Create, add, and return a new slide formatted with base conditions"""
        slide = self.slide_template.add_slide(presentation)
        self.professor_name_template(slide, presentation, professor_name)
        return slide


def build_room_slide_func(slide_template: SlideTemplate = SlideTemplate()) \
        -> Callable[[Presentation, str, str, Position, Dimension, str], Slide]:
    """Return a template which adds and returns a new slide formatted for a room"""
    base_slide_template = BaseSlideTemplate(slide_template)

    left, top = base_slide_template.professor_base_template.position
    width, height = base_slide_template.professor_base_template.dimensions

    room_defaults = TextboxTemplate(Position(left=left, top=top), Dimension(width=width, height=height))
    room_number_template = build_room_textbox_func(room_defaults)

    def add_room_slide(presentation: Presentation, professor_name: str, room_num: str,
                       room_pos: Position, room_dimensions: Dimension, floor_pic: str) -> Slide:
        """Create, add, and return a new slide formatted for a room flashcard"""
        slide = base_slide_template.add_base_slide(presentation, professor_name)
        room_number_template(slide, presentation, room_num)
        return slide

    return add_room_slide


def build_professor_slide_func(slide_template: SlideTemplate = SlideTemplate())\
        -> Callable[[Presentation, str, str], Slide]:
    """Return a template which adds and returns a new slide formatted for a professor picture"""
    base_slide_template = BaseSlideTemplate(slide_template)

    left, top = base_slide_template.professor_base_template.position
    width, height = base_slide_template.professor_base_template.dimensions

    room_defaults = TextboxTemplate(Position(left=left, top=top), Dimension(width=width, height=height))
    base_slide_template.room_number_template = build_room_textbox_func(room_defaults)

    def add_professor_slide(presentation: Presentation, professor_name: str, professor_pic: str) -> Slide:
        """Create, add, and return a new slide formatted for a professor picture flashcard"""
        slide = base_slide_template.add_base_slide(presentation, professor_name)
        picture = slide.shapes.add_picture(professor_pic, Inches(0), Inches(0))
        height_ratio = presentation.slide_height / picture.height
        picture.height = presentation.slide_height
        picture.width = int(picture.width * height_ratio)
        picture.left = (presentation.slide_width - picture.width) // 2
        return slide

    return add_professor_slide
