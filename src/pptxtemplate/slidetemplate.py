"""
Build slides as a template for use with python-pptx
"""
from pptx.presentation import Presentation
from pptx.slide import Slide

DEFAULT_BLANK_LAYOUT = 6


class SlideTemplate:
    """
    Build slides using this instances data as a template

    """
    def __init__(self, layout: int = DEFAULT_BLANK_LAYOUT):
        self.layout = layout

    def add_slide(self, presentation: Presentation) -> Slide:
        """Add a slide to presentation following the template and using these given details"""
        return presentation.slides.add_slide(presentation.slide_layouts[self.layout])
