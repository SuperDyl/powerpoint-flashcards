"""
Setup and create get-to-know-you flashcards using PowerPoint
"""

from pptxtemplate import SlideTemplate, TextboxTemplate, Position, Dimension, pptx_presentation, DEFAULT_BLANK_LAYOUT

from person import Person
from pptx import Presentation, slide as pptx_slide
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT

# FONT_COLOR_MODE = {"RGB": "rgb"}
BLANK_LAYOUT = DEFAULT_BLANK_LAYOUT


# class FlashcardBaseSlideTemplate(SlideTemplate):
#     """Base template for all Flashcard Slide Template classes"""
#     def __init__(self):
#         super().__init__()
#
#     def add_slide(self, presentation):
#         """Add a slide to presentation following the template and using these given details"""
#         slide = super().add_slide(presentation)
#         return slide
#
#
# class RoomSlideTemplate(FlashcardBaseSlideTemplate):
#     """Template for Professor room slides"""
#     def __init__(self):
#         super().__init__()
#
#     def add_slide(self, presentation, professor_name, room_number, floor_picture, room_dimensions, room_position):
#         """Add a slide to presentation following the template and using these given details"""
#         slide = super().add_slide(presentation)
#         textbox_height = self.name_bottom_padding
#         name_textbox = slide.shapes.add_textbox(left=0, top=presentation.height - textbox_height,
#                                                 width=presentation.width, height=textbox_height)
#         frame = name_textbox.text_frame
#         frame.text = professor_name
#         frame.auto_size = MSO_AUTO_SIZE.NONE
#         frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
#         font = frame.paragraphs[0].font
#         font.size = self.font_size
#         setattr(font.color, self.font_color_mode, self.font_color)
#         return slide

class PowerpointTextboxBaseTemplate:
    """
    Base formatting for all text boxes in PowerPoint Flashcards

    """
    def __init__(self, template: TextboxTemplate):
        self.template = template
        self.font_size = Pt(88)
        self.font_color = RGBColor(255, 0, 0)
        self.font_color_mode = "rgb"  # FONT_COLOR_MODE["RGB"]

    def add_base_textbox(self, slide: pptx_slide):
        """Add to slide and return a textbox with default formatting"""
        textbox = self.template.add_textbox(slide)

        frame = textbox.text_frame
        frame.auto_size = MSO_AUTO_SIZE.NONE
        frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

        font = frame.paragraphs[0].font
        font.size = self.font_size
        setattr(font.color, self.font_color_mode, self.font_color)

        return textbox


def professor_name_textbox_template(template: TextboxTemplate):
    """Return a function to format, add, and return a textbox formatted for a professor's name"""
    base_template = PowerpointTextboxBaseTemplate(template)

    def add_professor_name_textbox(slide: pptx_slide, professor_name: str):
        """Add to slide and then return a default formatted textbox"""
        textbox = base_template.add_base_textbox(slide)
        textbox.frame.text = professor_name
        return textbox

    return add_professor_name_textbox


def room_number_textbox_template(template: TextboxTemplate):
    """Return a function to format, add, and return a textbox formatted for a room number"""
    base_template = PowerpointTextboxBaseTemplate(template)

    def add_room_number_textbox(slide: pptx_slide, room_number: str):
        """Add to slide and then return a default formatted textbox"""
        textbox = base_template.add_base_textbox(slide)
        textbox.frame.text = room_number
        return textbox

    return add_room_number_textbox


class RoomSlideTemplate:
    """
    Add and return a new slide formatted for a room
    """
    def __init__(self, template: SlideTemplate):
        self.slide_template = template
        self.name_bottom_padding = Inches(1.58)

        left = Inches(0)
        top = presentation.height - textbox_height
        width = presentation.width
        height = textbox_height)

        self.room_number_template = RoomNumberTextboxTemplate(TextboxTemplate(Position(left=left, top=top), Dimension(width=width, height=height)))
        self.professor_name_template = ProfessorNameTextboxTemplate(TextboxTemplate(Position(left=left, top=top), Dimension(width=width, height=height)))

    def add_room_slide(self, presentation: pptx_presentation):
        """Create, add, and return a new slide formatted for a room flashcard"""
        slide = self.slide_template.add_slide(presentation)
        textbox_height = self.name_bottom_padding

        return slide


class FlashcardPowerPoint:
    def __init__(self, people: list = None):
        self.people = list(people) if people is not None else list()
        self.name_bottom_padding = Inches(1.58)
        self.font_size = Pt(88)
        self.font_color = RGBColor(255, 0, 0)
        self.font_color_mode = "rgb"

    def _add_default_slide(self, presentation: Presentation, person_name: str):
        slide = presentation.slides.add_slide(BLANK_LAYOUT)
        textbox_height = self.name_bottom_padding
        name_textbox = slide.shapes.add_textbox(left=0, top=presentation.height - textbox_height,
                                                width=presentation.width, height=textbox_height)
        frame = name_textbox.text_frame
        frame.text = person_name
        frame.auto_size = MSO_AUTO_SIZE.NONE
        frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        font = frame.paragraphs[0].font
        font.size = self.font_size
        setattr(font.color, self.font_color_mode, self.font_color)
        return slide

    def _add_room_slide(self, presentation: Presentation, person: Person, base_slide=None):
        if base_slide is None:
            base_slide = self._add_default_slide
        slide = base_slide(presentation, person.name)
        return slide

    def _add_picture_slide(self, presentation: Presentation, person: Person, base_slide=None):
        if base_slide is None:
            base_slide = self._add_default_slide
        slide = base_slide(presentation, person.name)
        return slide

    def _add_slides(self, presentation: Presentation, person: Person, room_slide=None, picture_slide=None):
        if room_slide is None:
            room_slide = self._add_room_slide
        if picture_slide is None:
            picture_slide = self._add_picture_slide

        if person.name and person.picture:
            room_slide(presentation, person)
        if person.name and person.room:
            picture_slide(presentation, person)

    def add_person(self, person: Person):
        self.people.append(person)

    def export_powerpoint(self, file_name):
        presentation = Presentation()
        for person in self.people:
            self._add_slides(presentation, person)

        presentation.save(file_name)
